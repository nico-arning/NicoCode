"""
This is a the python library of the BRIDGE LSC project written by Nicolas Arning.

Within this file you will find several classes that contain function for different action necessary
top process sequencing files and compute several analyses.

The functions are called within several short processing scripts, that are themselves called by
the nextflow pipeline bridge.nf

Almost all of the code within the dose response class is borrowed from Tommaso Mari with Permission

Any questions please email nicolas.arning@bayer.com
"""

import glob
import os
import pickle
import subprocess
import sys
import time
import warnings
from collections import defaultdict
from datetime import date
from functools import partial
from multiprocessing import Pool
from pathlib import Path
from zipfile import ZipFile

import anndata
import boto3
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import plotly.express as px
import scanpy as sc
import seaborn as sns
import tqdm
from matplotlib import rcParams
from matplotlib.backends.backend_pdf import PdfPages
from pptx import Presentation
from pptx.util import Cm, Pt
from scipy.io import mmread
from scipy.optimize import curve_fit, differential_evolution, minimize
from scipy.stats import f
from sklearn.decomposition import PCA
from sklearn.metrics import r2_score
from sklearn.preprocessing import MinMaxScaler
from umap import UMAP
from unipressed import IdMappingClient


class MyDefaultDict(dict):
    """
    Class for defaultdict to return key if not present
    """

    def __missing__(self, key):
        """
        Parameters
        ----------
        self : object
            Instance of the MyDefaultDict class

        key : string
            Dictionary key to be also used as value

        """
        return key + ": No annotation found"


class Datapond:
    """
    Class to interact with out S3 bucket behind the datapond
    """

    def __init__(
        self, local_path="./datapond", bucket_name="read-processing-bridge", cores=8
    ):
        """
        intialise the class with the bucket name and the local strucure

        Parameters
        ----------
        self : object
            Instance of the Datapond class

        local_path : string,
            Path of file to be uploaded

        bucket_name : string
            Name of the s3 bucket to be acessed

        cores : integer
            Number of cores to use

        """
        # Start boto client with credentials

        self._local_path = local_path
        self._bucket_name = bucket_name
        self._cores = cores

    def get_file_folders(self, filter_in_filename, prefix):
        """
        Get the file structure of the S3 bucket

        Parameters
        ----------
        self : object
            Instance of the Datapond class

        filter_in_filename : string
            Substring to filter paths for. If you want the Reads folder for example
            use the string 'Reads'.

        prefix : string
            The folder of the datapond to look in. Narrowing down folders
            significantly reduces processing time
        """

        file_names = []
        folders = []

        default_kwargs = {
            "Bucket": self._bucket_name,  # Look into this bucket
            "Prefix": prefix,  # Only look in this folder
        }
        next_token = ""

        while next_token is not None:  # Discover all folder under the prefix folder
            updated_kwargs = default_kwargs.copy()

            if next_token != "":
                updated_kwargs["ContinuationToken"] = next_token
            response = self._client.list_objects_v2(**default_kwargs)
            contents = response.get("Contents")

            if contents == None:
                return file_names, folders

            for result in contents:
                key = result.get("Key")
                if filter_in_filename in key:
                    if key[-1] == "/":
                        folders.append(key)
                    else:
                        file_names.append(key)
            next_token = response.get("NextContinuationToken")

        return file_names, folders

    def download_file_s3(self, file_name):
        """
        Download individual file from S3 bucket

        Parameters
        ----------
        self : object
            Instance of the Datapond class

        file_name : string
            Name of the file to be downloaded

        """

        local_path = Path(self._local_path)  # path of the folder destination
        file_path = Path.joinpath(local_path, file_name)
        file_path.parent.mkdir(parents=True, exist_ok=True)  # make parent folder

        self._client.download_file(self._bucket_name, file_name, str(file_path))

    def download_files(self, filter_in_filename, prefix):
        """
        Download files from S3 bucket

        Parameters
        ----------
        self : object
            Instance of the Datapond class

        filter_in_filename : string
            Substring to filter paths for. If you want the Reads folder for example
            use the string 'Reads'.

        prefix : string
            The folder of the datapond to look in. Narrowing down folders
            significantly reduces processing time

        """

        local_path = Path(self._local_path)

        # Get structure of the S3 bucket
        file_names, folders = self.get_file_folders(filter_in_filename, prefix)

        # Filter out only the paths that have the desired string

        if (
            filter_in_filename != ""
        ):  # If there is no filter then just download everything
            file_names = [x for x in file_names if filter_in_filename in x]

        for folder in folders:  # Loop through filesan make parents
            folder_path = Path.joinpath(local_path, folder)
            folder_path.mkdir(parents=True, exist_ok=True)

        file_names = [
            x for x in file_names if not Path.joinpath(local_path, x).is_file()
        ]

        # Loop over all files to download
        _ = list(
            tqdm.tqdm(map(self.download_file_s3, file_names), total=len(file_names))
        )

    def upload_files(self, filter_in_filename, prefix="./datapond/Experiments/"):
        """
        from current working directory, upload a 'local_dir' with all
        its subcontents (files and subdirectories...)
        to a aws bucket

        Parameters
        ----------

        self : object
            Instance of the Datapond class

        filter_in_filename : string
            Substring to filter paths for. If you want the Reads folder for example
            use the string 'Reads'.

        """

        # Use the following bash command to find the files to upload
        find_cmd = f"find {prefix} -type f -wholename *{filter_in_filename}*"

        # Run the bash command
        file_names = (
            subprocess.check_output(find_cmd, shell=True)
            .decode(sys.stdout.encoding)
            .split("\n")
        )

        # Loop over all files to upload
        _ = list(tqdm.tqdm(map(self.upload_file, file_names), total=len(file_names)))

    def upload_file(self, local_dir):
        """
        Upload files to s3 bucket

        Parameters
        ----------
        self : object
            Instance of the Datapond class

        local_dir : string
            Local directory to upload from

        """

        # Initialise client for accessing s3 bucket with credentials

        # The amazon path is the same as the local except the first three
        # folders
        aws_path = str("/".join(local_dir.split("/")[3:]))
        if aws_path != "":
            client.upload_file(local_dir, self._bucket_name, aws_path)


class DeSeqViz:
    """
    Class to visualise the output of DeSeq2. Requires the R part of the BRIDGE analysis first
    """

    def __init__(self, results_path, genome_dir="", alpha=0.1, cores=4, ntop_genes=500):
        """
        Parameters
        ----------
        self : object
            Instance of DeSeqViz class

        results_path : string
            Path of the Deseq results folder containing the tsv tables. Usually under <experiment folder>/DeSeq/Tables/

        meta_map : string
            Pickle file containing a dictionary with gene names as keys and annotation as values

        alpha : float
            Significance level. Usually 0.05

        cores : integer
            Number of cores to use

        ntop_genes : integer
            Number of top variance genes to use for PCA

        """

        self._results_path = results_path

        # Path to dump the output plots in
        self._plot_path = results_path.replace("Tables", "Plots")

        # Dictionary to style the 3D plots by
        self._style_grid = {
            "showgrid": True,
            "showticklabels": False,
            "showline": True,
            "linecolor": "silver",
            "linewidth": 4,
            "backgroundcolor": "rgba(0, 0, 0, 0.045)",
        }
        meta_map = glob.glob(genome_dir + "/eggnogg*.p")
        if len(meta_map) > 0:
            with open(meta_map[0], "rb") as meta_inf:
                self._meta_map = MyDefaultDict(pickle.load(meta_inf))
        else:
            self._meta_map = defaultdict(lambda: "No annotation found")

        self._alpha = alpha
        self._cores = cores
        self._ntop_genes = ntop_genes

    def well_pcas(self):
        """
        Wrapper function to make different PCA plots displaying the similarity
        between in gene expression between different wells

        Parameters
        ----------
        self : object
            Instance of DeSeqViz class

        """

    def read_in_tables4pca(input_tuble):
        """
        Process gene expression tables for PCA. Takes the output of
        the make_embedding function as input.

        Parameters
        ----------
        input_tuble : string
            Tuple containing meta data, 3d pca embedding and 2d embedding.
            This tuple gets created by the make_embedding function


        """

        meta_data, pca_embedding, pca_embedding_2d = input_tuble

        # Put all tables together
        meta_data = pd.concat(
            [meta_data, pca_embedding, pca_embedding_2d],
            axis=1,
        )

        # The following few lines give the control treatment
        # DMSO a random concentration assignment

        dmso_df = meta_data[meta_data["Treatments"] == "DMSO"]
        meta_data = meta_data[~(meta_data["Treatments"] == "DMSO")]
        concentrations = meta_data["Concentrations"].unique()
        meta_bin = [meta_data]

        for conc in concentrations:
            dmso_df["Concentrations"] = conc
            meta_bin.append(dmso_df.copy())
        meta_data = pd.concat(meta_bin, ignore_index=True)

        # Link concentration to fixed sizes of dots in the scatterplot
        logdoses = sorted(list(set(meta_data["Concentrations"].fillna(0))))
        size_dict = dict(zip(logdoses, range(1, len(logdoses) + 1)))
        meta_data["sizes"] = [
            (size_dict[x]) for x in meta_data["Concentrations"].fillna(0)
        ]

        # Make DMSO have a different shape
        meta_data["Symbol"] = "circle"
        meta_data.loc[meta_data["Treatments"] == "DMSO", "Symbol"] = "circle-open"
        meta_data.loc[meta_data["Treatments"] == "DMSO", "sizes"] = 1

        return meta_data

    def leiden_clustering(self):
        """
        Perform Leiden and Louvain clustering on the dataset.

        This method applies a series of preprocessing steps to the dataset,
        including filtering cells and genes, calculating quality control metrics,
        normalizing the data, identifying highly variable genes, performing
        principal component analysis (PCA), computing the neighborhood graph,
        and finally applying Leiden/Louvain clustering.

        Returns
        -------
        None

        """

        # Define parameters for preprocessing
        min_genes = 200  # Minimum number of genes expressed in a cell to be kept
        min_cells = 3  # Minimum number of cells expressing a gene to be kept
        target_sum = 1e4  # Target sum for normalization
        exclude_highly_expressed = True  # Whether to exclude highly expressed genes
        max_fraction = 0.2  # Maximum fraction for excluding highly expressed genes
        n_bins = 10  # Number of bins for scaling

        # Load the dataset
        X = pd.read_csv(
            f"{self._results_path}/../../full_raw_matrix_and_meta.csv", index_col=0
        )

        # Separate metadata columns
        meta = X.loc[:, X.columns.str.startswith("Metadata_")]
        meta.columns = meta.columns.str.replace("_ ", " ")
        meta.columns = meta.columns.str.replace("_", " ")
        if "Metadata Timepoints" in meta.columns:
            meta["Metadata Treatments"] = (
                meta["Metadata Treatments"]
                + " at "
                + meta["Metadata Timepoints"].astype(str)
                + "h"
            )

        # Keep only the non-metadata columns
        X = X.loc[:, ~X.columns.str.startswith("Metadata_")]

        # Fill missing values with the median of each column
        X = X.fillna(X.median())

        # Remove columns with zero variance and convert to integer type
        X = X.loc[:, X.var() != 0].astype(int)

        # Create an AnnData object
        adata = anndata.AnnData(X=X, obs=meta)
        
        print(meta)
        # Filter cells and genes based on the defined thresholds
        sc.pp.filter_cells(adata, min_genes=min_genes)
        sc.pp.filter_genes(adata, min_cells=min_cells)

        # Calculate quality control metrics
        sc.pp.calculate_qc_metrics(adata, percent_top=None, log1p=False, inplace=True)

        # Normalize the data
        sc.pp.normalize_total(
            adata,
            target_sum=target_sum,
            exclude_highly_expressed=exclude_highly_expressed,
            max_fraction=max_fraction,
            inplace=True,
        )
        adata.raw = adata

        # Log-transform the data
        sc.pp.log1p(adata)

        # Identify highly variable genes
        sc.pp.highly_variable_genes(adata)

        # Subset the data to include only highly variable genes
        adata = adata[:, adata.var.highly_variable]

        # Regress out the effects of total counts
        sc.pp.regress_out(adata, ["total_counts"])

        # Scale the data to unit variance and clip values exceeding the specified number of bins
        sc.pp.scale(adata, max_value=n_bins)

        # Perform Principal Component Analysis (PCA)
        sc.tl.pca(adata, svd_solver="arpack")

        # Compute the neighborhood graph
        sc.pp.neighbors(adata, n_neighbors=10, n_pcs=40)

        sc.tl.leiden(adata, resolution=1.0)

        # Run PAGA
        sc.tl.paga(adata, groups="leiden")
        sc.pl.paga(adata, plot=True)

        # Perform UMAP
        sc.tl.umap(adata, init_pos="paga")

        # Perform Leiden clustering
        sc.tl.leiden(adata)

        # Perform Louvain clustering
        sc.tl.louvain(adata)

        rcParams["figure.figsize"] = (14, 10)

        # Plot UMAP with LOUVAIN cluster labels and treatments
        sc.pl.umap(
            adata,
            color=["louvain", "Metadata Treatments"],
            title="UMAP plot with LOUVAIN cluster labels and treatments",
        )
        plt.savefig(f"{self._plot_path}/louvain_clustering_umap.png", dpi=300)
        plt.close()

        # Plot PCA with LOUVAIN cluster labels and treatments
        sc.pl.pca(
            adata,
            color=["louvain", "Metadata Treatments"],
            title="PCA plot with LOUVAIN cluster labels and treatments",
        )
        plt.savefig(f"{self._plot_path}/louvain_clustering_pca.png", dpi=300)
        plt.close()

        # Plot UMAP with LEIDEN cluster labels and treatments
        sc.pl.umap(
            adata,
            color=["leiden", "Metadata Treatments"],
            title="UMAP plot with LEIDEN cluster labels and treatments",
        )
        plt.savefig(f"{self._plot_path}/leiden_clustering_umap.png", dpi=300)
        plt.close()

        # Plot PCA with LEIDEN cluster labels and treatments
        sc.pl.pca(
            adata,
            color=["leiden", "Metadata Treatments"],
            title="PCA plot with LEIDEN cluster labels and treatments",
        )
        plt.savefig(f"{self._plot_path}/leiden_clustering_pca.png", dpi=300)
        plt.close()

        rcParams["figure.figsize"] = (14, 12)

        sc.tl.rank_genes_groups(adata, "leiden", method="wilcoxon")

        sc.pl.rank_genes_groups_matrixplot(
            adata,
            n_genes=4,
            values_to_plot="logfoldchanges",
            cmap="bwr",
            colorbar_title="log fold change",
        )
        plt.tight_layout()
        plt.subplots_adjust(bottom=0.5)  # Adjust the bottom margin as needed

        plt.savefig(f"{self._plot_path}/per_cluster_logfold.png", dpi=300)
        plt.close()

        sc.pl.rank_genes_groups_matrixplot(
            adata,
            n_genes=4,
            cmap="viridis",
            colorbar_title="Normalised gene expression",
        )
        plt.tight_layout()
        plt.subplots_adjust(bottom=0.5)  # Adjust the bottom margin as needed

        plt.savefig(f"{self._plot_path}/per_cluster_expression.png", dpi=300)
        plt.close()

        sc.pl.correlation_matrix(adata, "Metadata Treatments")
        plt.savefig(f"{self._plot_path}/correlation_treatments.png", dpi=300)
        plt.close()

        ax = sc.pl.correlation_matrix(adata, "leiden")
        plt.savefig(f"{self._plot_path}/correlation_clusters.png", dpi=300)
        plt.close()

        # Assuming 'Metadata_treatments' and 'louvain' are columns in your dataframe
        df = pd.DataFrame(
            {
                "Metadata Treatments": adata.obs["Metadata Treatments"],
                "leiden": adata.obs["leiden"],
                #'Known MoA': adata.obs["Metadata_Mode_of_Action"]
            }
        )
        df = df.sort_values("Metadata Treatments")

        cluster_treatment_counts = (
            df.groupby("leiden")["Metadata Treatments"]
            .value_counts()
            .unstack()
            .fillna(0)
            .T
        )
        cluster_treatment_counts.columns = [
            "Leiden cluster: " + x for x in cluster_treatment_counts.columns
        ]
        cluster_treatment_counts["Metadata Treatments"] = cluster_treatment_counts.index
        # cluster_treatment_counts["Metadata Treatments"] = [ x.split(" ")[0] for x in cluster_treatment_counts['Metadata Treatments']]
        cluster_treatment_counts.index = cluster_treatment_counts["Metadata Treatments"]
        data = cluster_treatment_counts.drop(["Metadata Treatments"], axis=1)
        mask = data == 0

        clustermap = sns.clustermap(
            cluster_treatment_counts.drop(["Metadata Treatments"], axis=1),
            figsize=(12, 12),
            cmap="Blues",
            annot=True,
            fmt="g",
            mask=mask,  # Apply the mask
        )
        plt.savefig(f"{self._plot_path}/clusters_membership.png", dpi=300)
        plt.close()

        return

    def replicate_correlations(self):
        """
        Calculate the correlations between the replicates

        Returns
        -------
        None

        """
        # Load the dataset
        df = pd.read_csv(
            f"{self._results_path}/../../full_corrected_matrix_and_meta.csv",
            index_col=0,
        )
        df.columns = df.columns.str.replace("_ ", "_")
        # Select the columns of interest including the metadata columns and the normalized counts
        metadata_cols = df.columns[df.columns.str.startswith("Metadata")]
        gene_cols = df.columns[~df.columns.str.startswith("Metadata_")]
        df.loc[ df["Metadata_Treatments"] == "DMSO", "Metadata_Concentrations"] = 0
        print(df)
        df["Replicate_col"] = df["Metadata_Treatments"] + "_" + df["Metadata_Concentrations"].astype(str)

        heatmap_dict = {}
        for replicate in df["Replicate_col"].unique():
            table = df.loc[df["Replicate_col"] == replicate, gene_cols]
            correlation_matrix = table.T.corr(method="pearson")
            # Identify the diagonal elements
            np.fill_diagonal(correlation_matrix.values, np.nan)
            # Set the diagonal elements to NA    # Calculate the average correlation
            average_column_correlation = correlation_matrix.mean().mean()
            # Print the average correlation
            heatmap_dict[replicate] = round(average_column_correlation, 3)

        # Extract keys and values from the dictionary
        labels = heatmap_dict.keys()
        values = heatmap_dict.values()

        labels, values = zip(*sorted(list(zip(labels, values)), key=lambda x: x[0]))
        # Create a bar plot
        # plt.hbar(values, labels)

        # Add labels and title
        fig = plt.figure(figsize=(5, 10))

        ax = sns.barplot(x=values, y=labels, palette="rainbow")
        for container in ax.containers:
            ax.bar_label(container, fontsize=10)

        plt.xlabel("Replicates")
        plt.ylabel("Correlation values normalised counts")
        #plt.title("Bar Plot of correlation between replicates")
        plt.tight_layout()

        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)

        plt.savefig(f"{self._plot_path}/replicate_correlations.png", dpi=300)

    def read_decay(self):
        """
        Get some basic statistics of the gene expression data like correlation
        in between treatments, then compute the number of genes covered at different
        read cut-offs


        Returns
        -------
        None

        """

        # Load the dataset
        X = pd.read_csv(
            f"{self._results_path}/../../full_raw_matrix_and_meta.csv", index_col=0
        )
        X.columns = X.columns.str.replace("_ ", " ")
        X.columns = X.columns.str.replace("_", " ")


        # Downsample to DMSO
        X = X[X["Metadata Treatments"] == "DMSO"]

        # Separate metadata columns
        meta = X.loc[:, X.columns.str.startswith("Metadata")]
        if "Metadata Timepoints" in meta.columns:
            meta["Metadata Treatments"] = (
                meta["Metadata Treatments"]
                + " at "
                + meta["Metadata Timepoints"].astype(str)
                + "h"
            )

        # Keep only the non-metadata columns
        X = X.loc[:, ~X.columns.str.startswith("Metadata")]

        # Fill missing values with the median of each column
        X = X.fillna(0).astype(int)

        # Create a plot to visualize the gene number at read cutoff
        plt.figure(figsize=(10, 10))

        # Define a read count cutoff
        read_count_cutoffs = [1, 10, 50, 100, 200, 500]  # Example cutoff values

        genes = []
        cutoffs = []
        datatypes = []
        X.index = range(X.shape[0])

        for i, row in X.iterrows():
            # Calculate the total read counts per gene for each dataframe

            # Count the number of genes that meet the cutoff for each dataframe
            gene_counts_at_cutoffs = [
                sum(row >= cutoff) for cutoff in read_count_cutoffs
            ]

            # Plot the gene number at read cutoff for each dataframe
            genes += gene_counts_at_cutoffs
            datatypes += [f"DMSO Replicate {i + 1}"] * len(read_count_cutoffs)
            cutoffs += read_count_cutoffs

        plt.figure(figsize=(5, 5))

        plot_df = pd.DataFrame(
            {
                "Genes covered at Read Cutoff": genes,
                "Read cut-off": cutoffs,
                "Sample": datatypes,
            }
        )

        ax = plt.gca()
        ax = sns.lineplot(
            data=plot_df,
            x="Read cut-off",
            y="Genes covered at Read Cutoff",
            hue="Sample",
            marker="o",
            palette="Set2",
            ax=ax,
        )
        plt.savefig(f"{self._plot_path}/read_decay_curves.png", dpi=300)

    def make_embedding(self):
        """
        Make PCA embedding of the wellplates based on the expression
        profiles of all genes.

        Parameters
        ----------
        self : object
            Instance of DeSeqViz class

        Returns
        -------
        meta_data : object
            Pandas dataframe containing meta data of the wellplate

        embedding3d : object
            Pandas dataframe containing 3D embedding.
            Rows are wells, columns are the Principal components 1 to 3

        embedding2d : object
            Pandas dataframe containing 2D embedding.
            Rows are wells, columns are the Principal components 1 and 2

        """

        # Read in corrected count data generated by DeSeq2
        data_frame = pd.read_csv(
            f"{self._results_path}/variance_stabilised_counts.csv", index_col=0
        ).T

        # Take only N top variance genes for PCA

        data_frame = data_frame.reindex(
            data_frame.var().sort_values(ascending=False).index, axis=1
        )
        data_frame = data_frame[data_frame.columns[: self._ntop_genes]]

        # Read in meta data
        meta_data = pd.read_csv(f"{self._results_path}/plate_matrix.csv", index_col=0)

        pca3d = PCA(n_components=3)
        pca2d = PCA(n_components=2)

        embedding3d = pd.DataFrame(
            pca3d.fit_transform(data_frame),
            index=data_frame.index,
            columns=["PC 1", "PC 2", "PC 3"],
        ).astype(float)
        embedding2d = pd.DataFrame(
            pca2d.fit_transform(data_frame),
            index=data_frame.index,
            columns=["2D PC 1", "2D PC 2"],
        ).astype(float)

        return meta_data, embedding3d, embedding2d

        def range4viewing(series):
            """
            Generate range for axis limits

            Parameters
            ----------
            series : object
                Pandas series of the column to be viewed

            Returns
            -------
            tuple : tuple
                tuple containing the minimum and maximum values for viewing scale

            """

            return (min(series) * 0.95, max(series) * 1.05)

        def slider_3d_plot(meta_data, cmap):
            """
            Make 3D plot of the PCA with a slider for viewing varying
            concentrations individually

            Parameters
            ----------
            meta_data : object
                Pandas dataframe containing meta data of the wellplate

            cmap : dictionary
                Dictionary containing mapping colours to treatments

            Returns
            -------
            fig: object
                Matplotlib figure containing the 3D PCA

            """

            fig = px.scatter_3d(
                meta_data,
                x="PC 1",
                y="PC 2",
                z="PC 3",
                color="Treatments",
                color_discrete_map=cmap,
                size="sizes",
                hover_data=["Concentrations", "Wellnames", "Barcodes", "Treatments"],
                opacity=0.6,
                labels={"Colour", "Sizes"},
                animation_frame="Concentrations",  # This does the slider magic
                range_x=range4viewing(meta_data["PC 1"]),
                range_y=range4viewing(meta_data["PC 2"]),
                range_z=range4viewing(meta_data["PC 3"]),
            )

            return fig

        def no_slider_3d_plot(meta_data, cmap):
            """
            Make 3D plot of the PCA without the slider

            Parameters
            ----------
            meta_data : object
                Pandas dataframe containing meta data of the wellplate

            cmap : dictionary
                Dictionary containing mapping colours to treatments

            Returns
            -------
            fig: object
                Matplotlib figure containing the 3D PCA

            """

            fig = px.scatter_3d(
                meta_data,
                x="PC 1",
                y="PC 2",
                z="PC 3",
                color="Treatments",
                color_discrete_map=cmap,
                size="sizes",
                hover_data=["Concentrations", "Wellnames", "Barcodes", "Treatments"],
                opacity=0.6,
                labels={"Colour", "Sizes"},
                range_x=range4viewing(meta_data["PC 1"]),
                range_y=range4viewing(meta_data["PC 2"]),
                range_z=range4viewing(meta_data["PC 3"]),
            )

            return fig

        def slider_2d_plot(meta_data, cmap):
            """
            Make 2D plot of the PCA with a slider for viewing varying
            concentrations individually

            Parameters
            ----------
            meta_data : object
                Pandas dataframe containing meta data of the wellplate

            cmap : dictionary
                Dictionary containing mapping colours to treatments

            Returns
            -------
            fig: object
                Matplotlib figure containing the 2D PCA

            """

            # Make groups that identify the same treatment over different
            # concentrations
            meta_data = meta_data.fillna("None_")
            meta_data_animation_group = [
                f"{x[0]}|{x[1]}"
                for x in zip(meta_data["Wellnames"], meta_data["Treatments"])
            ]
            meta_data_animation_group = [
                f"""{x.split("|")[0].split("_")[0]}{x.split("_")[1]}"""
                for x in meta_data_animation_group
            ]
            meta_data["meta_data_animation_group"] = meta_data_animation_group

            px.scatter(
                meta_data,
                x="2D PC 1",
                y="2D PC 2",
                color="Treatments",
                color_discrete_map=cmap,
                size="sizes",
                hover_data=["Concentrations", "Wellnames", "Barcodes", "Treatments"],
                opacity=0.6,
                labels={"Colour", "Sizes"},
                animation_frame="Concentrations",
                animation_group="meta_data_animation_group",
            )

            return fig

        def normal_2d_plot(meta_data):
            """
            Make 2D plot of the PCA without a slider

            Parameters
            ----------
            meta_data : object
                Pandas dataframe containing meta data of the wellplate

            Returns
            -------
            fig: object
                Matplotlib figure containing the 2D PCA

            """
            fig = plt.figure()
            sns.scatterplot(
                meta_data,
                x="2D PC 1",
                y="2D PC 2",
                hue="Treatments",
                style="Concentrations",
                alpha=0.8,
            )
            return fig

        def make_pca(self, meta_data, cmap, type_pca):
            """
            Wrapper around the plotting functions. Make different types
            of PCA plots

            Parameters
            ----------
            self : object
                Instance of DeSeqViz class

            meta_data :
                Pickle file containing a dictionary with gene names as keys and annotation as values

            cmap : object
                Pandas dataframe containing meta data of the wellplate

            type_pca : string
                Type of plot to generate

            Returns
            -------
            fig: object
                Matplotlib figure containing the PCA plots

            """

            # Find out which plot to make, then make it
            if type_pca == "3d_slider":
                fig = slider_3d_plot(meta_data, cmap)
                three_d = True
            elif type_pca == "3d_no_slider":
                fig = no_slider_3d_plot(meta_data, cmap)
                three_d = True
            else:
                fig = slider_2d_plot(meta_data, cmap)
                three_d = False

            # Give the plot a cube appearance
            fig.update_layout(scene_aspectmode="cube")
            scene_layout = {
                "showgrid": True,
                "showticklabels": False,
                "showline": True,
                "linecolor": "silver",
                "linewidth": 4,
                "backgroundcolor": "rgba(0, 0, 0, 0.045)",
            }

            # Some adjustments
            if three_d:
                fig.update_layout(
                    scene={
                        "xaxis": scene_layout,
                        "yaxis": scene_layout,
                        "zaxis": scene_layout,
                    }
                )
            else:
                fig.update_layout(scene={"xaxis": scene_layout, "yaxis": scene_layout})

            return fig

        # Read in files
        meta_data = read_in_tables4pca(make_embedding(self))

        # Make colourmap
        cmap = sns.color_palette(
            "tab20", n_colors=meta_data["Treatments"].nunique()
        ).as_hex()
        cmap = dict(zip(meta_data["Treatments"].unique(), cmap))
        cmap["DMSO"] = "#454545"

        meta_data["Concentrations"] = meta_data["Concentrations"].astype(float)
        meta_data = meta_data.sort_values(by=["Concentrations"])
        treatment_per_concentration = meta_data[
            meta_data["Treatments"] != "DMSO"
        ].groupby(["Concentrations"])["Treatments"]
        n_treatment_per_concentration = treatment_per_concentration.nunique()

        # Find concentration with only one treatment
        conc_only_one_treat = n_treatment_per_concentration[
            n_treatment_per_concentration <= 1
        ]
        for conc in conc_only_one_treat.index:
            print(f"Concentration {conc} only has one treatment")
            meta_data = meta_data[meta_data["Concentrations"] != conc]

        fig = make_pca(self, meta_data, cmap, "3d_slider")
        fig.write_html(
            f"{self._plot_path}/3D_umap_slider.html",
            auto_play=False,
        )
        fig = make_pca(self, meta_data, cmap, "3d_no_slider")
        fig.write_html(
            f"{self._plot_path}/3D_umap_no_slider.html",
            auto_play=False,
        )

        fig = make_pca(self, meta_data, cmap, "2d_slider")
        fig.write_html(
            f"{self._plot_path}/2D_pca_slider.html",
            auto_play=False,
        )

        fig = normal_2d_plot(meta_data)
        fig.savefig(f"{self._plot_path}/2D_umap.png", dpi=300)

    def volcano_plot(self):
        """
        Make volcano plots per treatment based on the differential expression profiles

        Parameters
        ----------
        self : object
            Instance of DeSeqViz class

        """

        def make_gene_atlas(results_path, cores, pca_first=False):
            """
            Make a UMAP embedding of the difference in expression
            of individual genes. Here the data from all different treatments are used.
            Later by treatment change in expression are used for highlighting.
            In the resulting embedding every row is a gene and the 2 columns.
            are the embedding

            Parameters
            ----------
            results_path : string
                Path of the Deseq results folder containing the tsv tables. Usually under <experiment folder>/DeSeq/Tables/

            cores : integer
                Number of cores to use

            pca_first : bool
                Whether to do a PCA first as a preprocessing step

            Returns
            -------
            umap_embedding : object
                Pandas dataframe containing Umap embedding. Rows are genes. Columns is the UMAP embedding

            """

            # Read the differential expression data
            data_frame = pd.read_csv(
                f"{results_path}/variance_stabilised_counts.csv", index_col=0
            )
            backup_index = data_frame.index

            # First use a PCA for dimensionality reduction
            if pca_first:
                pca = PCA(n_components=10)
                data_frame = pca.fit_transform(data_frame)

            # Initialise Umap
            umap = UMAP(
                n_components=2,
                verbose=True,
                n_jobs=cores,
                init="random",
                metric="euclidean",
                n_neighbors=8,
                min_dist=0.001,
                learning_rate=0.1,
                random_state=42,
            )

            # Transform data
            umap_embedding = pd.DataFrame(
                umap.fit_transform(data_frame), columns=["Umap 1", "Umap 2"]
            )

            # Reindex with original index
            umap_embedding.index = backup_index

            return umap_embedding

        def plot_gene_atlas(table, umap_embedding, plot_path, scaler):
            """
            Make 2D scatter plot of the UMAP embedding of differential expression.
            Colouring is based on log 2 fold change of gene expression as calculated
            by Deseq. Every dot represents one gene.

            Parameters
            ----------
            table : string
                Path of the per treatment Deseq output table

            umap_embedding : object
                Pandas dataframe containing Umap embedding. Rows are genes. Columns is the UMAP embedding

            plot_path : string
                Path to put plots in.

            scaler : object
                Min max scaler object to use, so that the log2foldchange will have the same colouring across
                all treatments

            """

            # Get the name of the treatment out of the table name.
            treatment = table.split("/")[-1].split("_")[0]

            data_frame = pd.read_table(table, index_col=0)

            # Get the data into the correct format
            data_frame["-Log10p"] = -1 * np.log10(data_frame["padj"])
            data_frame["Significance"] = "Not significant"
            data_frame["Type hit"] = "Below Threshold"
            data_frame["symbols"] = 1
            data_frame["padj"] = data_frame["padj"].astype(float)

            # Highlight interesting genes. As in over log2foldchange 1 and over
            # significance of alpha
            data_frame.loc[data_frame["log2FoldChange"].abs() > 1, "Type hit"] = (
                "Absolute log fold change over 1"
            )
            data_frame.loc[
                (data_frame["log2FoldChange"].abs() > 1)
                & (data_frame["padj"] < self._alpha),
                "Type hit",
            ] = """Absolute log fold change over 1
                   \n and hit over significance threshold of 0.05"""

            # Scale log2foldchange
            data_frame["log2FoldChange"] = scaler.transform(
                data_frame["log2FoldChange"].values.reshape(-1, 1)
            )

            # Merge with the embedding dataframe
            data_frame = data_frame.merge(
                umap_embedding, left_index=True, right_index=True
            )

            # Make 3D plot
            fig = plt.figure(figsize=(10, 10))
            axis = plt.gca()
            cmap = sns.color_palette("coolwarm", as_cmap=True)

            # Plot uninteresting genes first
            sns.scatterplot(
                data_frame[data_frame["Type hit"] == "Below Threshold"],
                x="Umap 1",
                y="Umap 2",
                c="silver",
                size=1,
                alpha=0.1,
                ax=axis,
                legend=False,
            )

            # Then plot the highglighted genes over
            sns.scatterplot(
                data_frame[data_frame["Type hit"] != "Below Threshold"],
                x="Umap 1",
                y="Umap 2",
                hue="log2FoldChange",
                palette=cmap,
                alpha=0.9,
                size=8,
                ax=axis,
            )

            fig = plt.gcf()
            plt.tight_layout()
            fig.savefig(f"{plot_path}/{treatment}_gene_atlas.png", dpi=300)

        def translate_gene_id2uniprotkb(data):
            """
            Function to translate gene_ids to uniprotkb
            """

            ids = list(data.index)
            if ids[0].startswith("Bcin"):
                ids = [x.replace("Bcin", "BCIN_") for x in ids]

            try:
                request = IdMappingClient.submit(
                    source="Gene_Name", dest="UniProtKB", ids=ids
                )
            except KeyError:
                gene_id2hyperlink = MyDefaultDict({})
                return gene_id2hyperlink

            result = None
            while result is None:
                try:
                    # connect
                    result = list(request.each_result())
                except:
                    time.sleep(3.0)
                    pass
            result = [x["to"] for x in result]
            hyperlink = [f"https://www.uniprot.org/uniprotkb/{x}/entry" for x in result]
            gene_id2hyperlink = dict(zip(data.index, hyperlink))
            gene_id2hyperlink = MyDefaultDict(gene_id2hyperlink)
            return gene_id2hyperlink

        def make_volcano(self, table, treatment):
            """
            Make an interactive volcano plot as an html. Attach meta data to individual genes.

            Parameters
            ----------
            self : object
                Instance of DeSeqViz class

            table : string
                Path of the per treatment Deseq output table

            treatment : string
                Name of the treatment

            Returns
            -------
            fig : object
                The figure object of the volcano plot

            pd.Series : object
                Pandas series containing the log2foldchange. Will be used later for
                colouring in the plot_gene_atlas function

            """

            # Read in data
            data_frame = pd.read_table(table, index_col=0)

            # Convert p-values into -Log10p
            data_frame["-Log10P"] = np.log10(data_frame["padj"]) * -1
            data_frame.loc[data_frame["-Log10P"] < 0, "-Log10P"] = 0

            # Divide genes/rows into over and under log2folhange and
            # signficiance cut-offs
            data_frame["Type hit"] = "Under threshold"
            data_frame.loc[data_frame["log2FoldChange"] > 1, "Type hit"] = (
                "Log 2 fold change over 1"
            )
            data_frame.loc[data_frame["padj"] < self._alpha, "Type hit"] = (
                f"P-value smaller than {self._alpha}"
            )
            data_frame.loc[
                (data_frame["padj"] < self._alpha) & (data_frame["log2FoldChange"] > 1),
                "Type hit",
            ] = f"Log 2 fold change over 1 and p-value over {self._alpha}"
            data_frame.loc[
                (data_frame["padj"] < self._alpha)
                & (data_frame["log2FoldChange"] < -1),
                "Type hit",
            ] = f"Log 2 fold change over 1 and p-value over {self._alpha}"
            data_frame.loc[(data_frame["padj"] > self._alpha), "meta"] = (
                "Under Threshold"
            )
            data_frame.loc[
                (data_frame["log2FoldChange"] < 1)
                & (data_frame["log2FoldChange"] > -1),
                "meta",
            ] = "Under Threshold"
            data_frame = data_frame.sort_values("padj")
            hyperlinks_dict = translate_gene_id2uniprotkb(data_frame.iloc[:10])

            top_genes = data_frame.iloc[:15]
            uniprot_links = list(top_genes.index)
            uniprot_links = [f"{x}\t{hyperlinks_dict[x]}" for x in uniprot_links]

            fig = plt.figure(figsize=(10, 10))

            # Make figure
            g = sns.scatterplot(
                data=data_frame,
                x="log2FoldChange",
                y="-Log10P",
                hue="Type hit",
                alpha=0.8,
                palette=["silver", "cornflowerblue", "mediumseagreen", "crimson"],
                legend=False,
            )
            plt.title = (f"Volcano plot of treatment {treatment}",)

            # Add lines to show thresholds
            ax = plt.gca()
            # Add lines to show thresholds
            ax.axhline(y=np.log10(0.05) * -1, lw=2, ls="--", c="silver")

            ax.axvline(x=1, lw=2, ls="--", c="silver")

            ax.axvline(x=-1, lw=2, ls="--", c="silver")

            texts = []

            i = 1
            for gene, row in top_genes.iterrows():
                texts.append(
                    ax.text(
                        row["log2FoldChange"],
                        row["-Log10P"],
                        f"{i}: {gene}",
                        fontsize=8,
                    )
                )

                i += 1
                # print(texts)
                # adjust_text(
                #    texts,
                #    only_move={'points':'y', 'texts':'y', "explode": "xy", "text": "xy", "static": "xy", "pull": "xy"},
                #    arrowprops=dict(arrowstyle="-", color='grey', ls = "--", lw=0.5,))

            return fig, data_frame["log2FoldChange"], uniprot_links

        def volcano_per_treatment(table, plot_path):
            """
            Wrapper around make_volcano so we can map the funcction more seasily

            Parameters
            ----------

            table : string
                Path of the per treatment Deseq output table

            plot_path : string
                Path to put plots in.

            Returns
            -------
            pd.Series : object
                Pandas series containing the log2foldchange. Will be used later for
                colouring in the plot_gene_atlas function

            """

            # Get treatment name through the table
            treatment = table.split("/")[-1].split("_")[0]

            # Make volcano plot then save it
            fig, log2foldchange, uniprot_links = make_volcano(self, table, treatment)
            with open(f"{plot_path}/{treatment}_uniprot_links.txt", "w+") as outf:
                outf.write("\n".join(uniprot_links))
            fig.savefig(f"{plot_path}/{treatment}_volcano_plot.png")
            return log2foldchange

        # Path of the DeSeq output
        results_path = self._results_path

        # Dictionary containing meta data per gene
        meta_map = self._meta_map

        # All differential expression tables per treatment
        output_tables = list(glob.glob(f"{results_path}/*_differential_expression.tsv"))

        # Apply the making of the volcano plot over the list of treatment tables
        # Then extract the log2foldchange data per treatment
        func = partial(volcano_per_treatment, plot_path=self._plot_path)
        log2foldchanges = list(
            tqdm.tqdm(map(func, output_tables), total=len(output_tables))
        )

        # Put the log2dfoldchanges into one data table
        log2foldchanges = pd.concat(log2foldchanges, ignore_index=True)

        # Scale the log2foldchanges for steady colouring across all treatments
        # scale to max log2foldchange
        scale_factor = max(log2foldchanges.abs())
        scaler = MinMaxScaler().fit(np.array([[scale_factor * -1], [scale_factor]]))

        # Make gene atlas - Represeantation of differential expression on gene
        # level
        # gene_atlas_embedding = make_gene_atlas(self._results_path, self._cores)
        # func = partial(
        #    plot_gene_atlas,
        #    plot_path=self._plot_path,
        #    umap_embedding=gene_atlas_embedding,
        #    scaler=scaler,
        # )
        # _ = list(tqdm.tqdm(map(func, output_tables), total=len(output_tables)))


class DoseResponse:
    """
    Class to visualise dose response curves from deseq results.
    A large part of this class is copied with persmission from Tommaso Mari

    """

    def __init__(self, results_path, genome_dir="", cores=4, alpha=0.05):
        """
        Initialisation of DoseRespons<e class

        Parameters
        ----------
        self : object
            Instance of DoseResponse class

        results_path : string
                    Path of the Deseq results folder containing the tsv tables. Usually under <experiment folder>/DeSeq/Tables/

        alpha : float
            Significance level. Usually 0.05

        cores : integer
            Number of cores to use

        meta_map : string
            Pickle file containing a dictionary with gene names as keys and annotation as values
        """

        self._results_path = results_path
        self._plot_path = results_path.replace("Tables", "Plots")
        self._cores = cores
        self._alpha = alpha

        meta_map = glob.glob(genome_dir + "/eggnogg*.p")
        if len(meta_map) > 0:
            with open(meta_map[0], "rb") as meta_inf:
                self._meta_map = MyDefaultDict(pickle.load(meta_inf))
        else:
            self._meta_map = defaultdict(lambda: "No annotation found")

    def model_sigmoids(
        self,
        row,
        dmso_data,
        xvals,
        meta_map,
        median=False,
        n_param=4,
        mode="min",
        bottom_bounds=[-0.5, 0.5],
    ):
        """Main wrapper to calculate sigmoids

        Parameters
            ----------
            row : pd.Series
                Experimental values to use (your y-axis).
                This can be also a row of a dataframe, if the function is applyed

            xvals : pd.Series
                Values on the x-axis to calculate sigmoids (e.g. Cmpd concentrations).
                Must be of the same length as row

            median : bool, list
                if False, all values of row are used for calculation.
                If list: the values indicate the replicates in your experiment
                one unique label per experiment
                (e.g. ['conc1', 'conc1', 'conc2', 'conc2', ...])
                if specified, the sigmoid will be calculated using the median values
                for each replicate.
                R2 and F-statistic will still be calculated using the full data

            plot : bool
                if True, generates one plot for each sigmoids.
                As of now, just outputs to console.

            n_param : integer
                either 3 or 4. Specifies which sigmoid function to use.
                With 3, the left of the sigmoid is fixed at 1
                With 4, the left of the sigmoid can be are specified by bottom_bounds.
            mode : string
                Either "min" or "fit"
                Specifies the use of two different ways of calculating the sigmoid.
                NOTE: "min" can only be used with n_param=4
                "min" is the recommended way

            bottom_bounds : list
                list of length 2, contains bottom and top values for the left side
                of the sigmoid.
                If using log-transformed values for row (hence start of the curve is
                at y=0), suggested is [-.5, .5] (default)
                If row non log transformed (hence start of the curve is at y=1),
                suggested is [.5, 1.5]
                This range **must** be symmetric to the expected start of the curve

            Returns
            -------

            pd.Series
                Curve parameters and statistics
                    Bottom (optional)  : left-most value of y of the curve
                        only present if n_param = 4
                    Top : right-most value of y of the curve
                    EC50 : point of flex (f"(x)=0) of the sigmoid
                    Slope : slope of the sigmoid
                    R2 : residual of sum squared between observed data
                        and calculated sigmoid
                    p-value : result of F-statistic calculation. The test tells
                        the likelyhood of the data to originate from a straight line
                        rather than the calculated sigmoid
        """

        def sigmoid_4p(x, b, t, e, s):
            """
            Compute a four parameter sigmoid curve

            Returns
            -------
            float
                y value for the given parameters

            """

            return b + (t - b) / (1 + np.exp(s * (np.log(x) - np.log(e))))

        def sigmoid_3p(x, b, e, s):
            """
            Compute a three parameter sigmoid curve

            Returns
            -------
            float
                y value for the given parameters


            """

            return b + (1 - b) / (1 + np.exp(s * (np.log(x) - np.log(e))))

        def generate_initial_parameters(
            data_frame, max_conc, bottom_bounds=[-0.5, 0.5]
        ):
            """
            Create initial parameters for the sigmoid using the differential
            evolution method to find the global minimum of sum of
            squared error for every parameter of the multivariate
            sigmoid function

            Parameters
            ----------
            data_frame : object
                Pandas dataframe containing the concentrations and gene counts

            max_conc : float
                maximum concentration

            bottom_bounds : list
                bottom bounds for the sigmoid (left side of the curve)

            """

            max_conc *= 1.1
            parameterBounds = []
            parameterBounds.append([-100, 100])  # search bounds for top
            parameterBounds.append(bottom_bounds)  # search bounds for bottom
            parameterBounds.append([0.0, max_conc])  # search bounds for flex
            parameterBounds.append([-500, 500])  # search bounds for slope

            def sumOfSquaredError(
                parameterTuple, x_data=data_frame.x, y_data=data_frame.y
            ):
                """
                Generate the sum of squares error.

                Parameters
                ----------
                parameterTuple : tuple
                    Tuple containing parameters

                x_data : object
                    Pandas series containing the x values

                y_data : object
                    Pandas series containing the y values

                Returns
                -------
                float
                    The sum of squares error

                """
                # do not print warnings by genetic algorithm
                warnings.filterwarnings("ignore")
                val = sigmoid_4p(x_data, *parameterTuple)
                return np.sum((y_data - val) ** 2.0)

            result = differential_evolution(sumOfSquaredError, parameterBounds, seed=3)
            return result.x

        def curve_fit_values(data_frame, p_number=4):
            """
            Fit the sigmoid curve function

            Parameters
            ----------
            data_frame : object
                Pandas dataframe containing the concentrations and gene counts

            p_number : integer
                Number of parameters for the sigmoid function

            Returns
            -------
            object
                Fitted sigmoid curve


            """

            # Check whether to use 3 or 4 parameter sigmoid curve
            if p_number == 4:
                fun = sigmoid_4p
                empty = (np.nan, np.nan, np.nan, np.nan)
            elif p_number == 3:
                fun = sigmoid_3p
                empty = (np.nan, np.nan, np.nan)

            try:
                fitted = curve_fit(
                    f=fun,
                    xdata=data_frame.x,
                    ydata=data_frame.y,
                )[0]
            except RuntimeError:  # If the sigmoid can't be fitted return NAs
                fitted = empty

            return fitted

        def curve_fit_minimize(data_frame, max_conc, bottom_bounds=[-0.5, 0.5]):
            """
            Fit curve by Nelder Mead method


            Parameters
            ----------
            data_frame : object
                Pandas dataframe containing the concentrations and gene counts

            max_conc : float
                maximum concentration

            bottom_bounds : list
                bottom bounds for the sigmoid (left side of the curve)

            Returns
            -------
            object
                solution array

            """

            # Get initial parameters by differentioal evolution method
            p0 = generate_initial_parameters(data_frame, max_conc, bottom_bounds)

            # Minimise function
            fitted = minimize(
                driver_func,
                p0,
                args=(data_frame.x, data_frame.y),
                bounds=((None, None), tuple(bottom_bounds), (0, np.Inf), (None, None)),
                method="Nelder-Mead",
            )
            return fitted.x

        def driver_func(x, xobs, yobs):
            """
            Driver function for the minimisation

            Parameters
            ----------
            x : object
                Pandas dataframe containing the concentrations and gene counts

            xobs : float
                maximum concentration

            yobs : list
                bottom bounds for the sigmoid (left side of the curve)

            Returns
            -------
            object
                solution array

            """
            # Evaluate the fit function with the current parameter estimates

            ynew = sigmoid_4p(xobs, *x)
            yerr = RSS(yobs, ynew)

            return yerr

        def RSS(y, y_pred):
            """
            Compute residual sum of squares error

            Parameters
            ----------
            y : float
                real y

            y_pred : float
                predicted y

            Returns
            ----------
            float
                Residual sum of squares error

            """
            return np.sum(np.square(y_pred - y))

        # Read in a couple of stats from the input
        gene_name, row, dmso_df, padj, logfoldchange, treatment = row

        # Create a dataframe exclusively for plotting as a background
        # Behind the sigmoid curve to be fitted

        plot_df = pd.DataFrame(
            {
                "Corrected counts": pd.concat(
                    [row, dmso_df[gene_name]], ignore_index=True
                ),
                "Concentrations": pd.concat(
                    [xvals, dmso_df["Concentrations"]], ignore_index=True
                ),
                "Treatments": [treatment] * xvals.shape[0]
                + ["DMSO"] * dmso_df.shape[0],
            }
        )

        # Scale row so its bound to 0 and 1
        row = (row - row.min()) / (row.max() - row.min())
        row_original = row
        xvals_original = xvals

        # Do the same for the plotting
        plot_df["Corrected counts"] = (
            plot_df["Corrected counts"] - plot_df["Corrected counts"].min()
        ) / (plot_df["Corrected counts"].max() - plot_df["Corrected counts"].min())

        # Get the sigmoid function
        sigm_fun = sigmoid_3p
        if n_param == 4:
            sigm_fun = sigmoid_4p

        xvals.index = row.index

        # Make dataframe for fitting
        df_bu = pd.DataFrame({"x": xvals, "y": row}).dropna()

        # If method median fit to median of corrected counts
        if median:
            new_ind = row.index.to_series()
            new_ind = new_ind.str.split("_").str[:-1].apply(lambda y: "_".join(y))

            row = pd.Series(np.unique(xvals)).apply(
                lambda y: np.nanmedian(row.loc[xvals == y])
            )
            xvals = np.unique(xvals)

        # Create dataframe out of input data

        data_frame = pd.DataFrame({"x": xvals, "y": row})
        data_frame = data_frame.dropna()

        # Fit the sigmoid
        if mode == "fit":
            out = curve_fit_values(data_frame, p_number=n_param)
        elif mode == "min":
            out = curve_fit_minimize(data_frame, data_frame.x.max(), bottom_bounds)

        # Create dataframe containing details about the fitted sigmoid
        if n_param == 3:
            out_ind = ["Top", "EC50", "Slope"]

        elif n_param == 4:
            out_ind = ["Bottom", "Top", "EC50", "Slope"]
        out = pd.Series(out, index=out_ind)

        # Initialise values
        r2 = np.nan
        pv_F = np.nan
        if np.all(out.isna()):
            r2 = np.nan
        try:
            r2 = r2_score(df_bu.y, sigm_fun(df_bu.x, *out))
        except ValueError:
            r2 = np.nan

        # Compute different scorers
        my_zero_y = np.sum(bottom_bounds) / 2
        rss_sig = RSS(df_bu.y, sigm_fun(df_bu.x, *out))
        rss_lin = RSS(df_bu.y, np.repeat(my_zero_y, np.size(df_bu.x)))
        d1 = np.size(out) - 1
        d2 = np.size(df_bu.y) - np.size(out)
        f_stat = ((rss_lin - rss_sig) / rss_sig) * (d2 / d1)
        pv_F = f.sf(f_stat, dfn=d2, dfd=d1)
        if np.all(out.isna()):
            pv_F = np.nan

        # Plot the dose response raw data as a backdrop for the fitted sigmoid
        fig, axes = plt.subplots(1, 2, figsize=(20, 10))
        sns.scatterplot(
            plot_df[plot_df["Treatments"] == treatment],
            x="Concentrations",
            y="Corrected counts",
            alpha=0.8,
            ax=axes[0],
        )

        # Get the DMSO data
        plot_df = plot_df[plot_df["Treatments"] == "DMSO"]
        high_end_dmso = plot_df["Corrected counts"].max()
        low_end_dmso = plot_df["Corrected counts"].min()
        median_dmso = plot_df["Corrected counts"].median()

        # Put in lines for the DMSO data
        axes[0].axhline(high_end_dmso, ls="--", c="silver", alpha=0.7, lw=2)
        axes[0].axhline(low_end_dmso, ls="--", c="silver", alpha=0.7, lw=2)
        axes[0].axhline(median_dmso, c="silver", alpha=1, lw=2)
        axes[0].text(
            max(data_frame.x),
            high_end_dmso,
            "High end DMSO",
            size="medium",
            c="dimgrey",
            ha="right",
        )
        axes[0].text(
            max(data_frame.x),
            low_end_dmso,
            "Low end DMSO",
            size="medium",
            c="dimgrey",
            ha="right",
        )
        axes[0].text(
            max(data_frame.x),
            median_dmso,
            "Median DMSO",
            size="medium",
            c="dimgrey",
            ha="right",
        )

        # Compute points to show fitted sigmoid curves
        predx = np.linspace(min(data_frame.x) * 0.9, max(data_frame.x) * 1.1, 100)
        predy = sigm_fun(predx, *out)

        out["R2"] = r2
        out = out.round(2)

        # Plot fitted sigmoid
        sns.lineplot(
            x=predx,
            y=predy,
            ax=axes[0],
            alpha=0.8,
            label="Fitted Sigmoid",
            lw=2,
            dashes=True,
            legend=True,
        )

        axes[0].set(xscale="log")
        props = {"boxstyle": "round", "facecolor": "lightgray", "alpha": 0.6}

        # Make an annotation box with data about the sigmoid fit
        textstr = f"""Gene:\t{gene_name}\np-value diff. expr.:
        \t{padj:.3e}\nLog2-fold change expr.\t{logfoldchange:.2f}\n\nFitted Sigmoid\nSlope:\t{out["Slope"]}
        \nR2:\t{out["R2"]}\np-value:\t{pv_F:.3e}
        \nEC50:\t{out["EC50"]}""".expandtabs()
        axes[0].text(
            0.05,
            0.05,
            textstr,
            transform=axes[0].transAxes,
            fontsize=12,
            verticalalignment="bottom",
            bbox=props,
        )
        axes[0].tick_params(axis="x", which="minor")
        axes[0].set_ylabel("Scaled corrected counts")

        meta_of_interest = [
            "Description",
            "Preferred_name",
            "GOs",
            "PFAMs" "KEGG_Pathway",
            "BRITE",
        ]
        # Plot metadata of the gene alongside dose response curves
        axes[1].set_axis_off()
        gene_name = gene_name.strip()
        axes[1].text(
            0,
            axes[1].get_ylim()[1],
            meta_map[gene_name],
            bbox=props,
            size="medium",
            ha="left",
            va="top",
        )

        axes[0].set_title(label="Dose response with fitted sigmoid", size="large")
        axes[1].set_title(label="Gene annotation", size="large")
        out["R2"] = r2
        out["p-value"] = pv_F

        if out["Slope"] < 0:
            direction = ["Negative"] * len(xvals_original)
        else:
            direction = ["Positive"] * len(xvals_original)

        return (
            row_original,
            pd.Series(xvals_original),
            pd.Series(direction),
            r2,
            gene_name,
            out,
            fig,
        )

    def dose_response(self):
        """
        Wrapper function around the model_sigmoids function supplied
        by tommaso mari. Wrote this so I could easily apply
        this to a lot of different genes in a highly parallelised
        setting

        Parameters
        ----------
        self : object
            Instance of DoseResponse class

        """

        def process_slope_df(slope_df, alpha, table):
            slope_df["Gene"] = slope_df.index
            slope_df["EC50_rank"] = slope_df["EC50"].rank(ascending=True)
            slope_df["R2_rank"] = slope_df["R2"].rank(ascending=False)
            slope_df["p-value_rank"] = slope_df["padj"].rank(ascending=True)
            slope_df["Foldchange_rank"] = (
                slope_df["log2FoldChange"].abs().rank(ascending=False)
            )
            slope_df["Combined_rank"] = (
                slope_df[["R2_rank", "p-value_rank", "Foldchange_rank", "EC50_rank"]]
                .mean(axis=1)
                .rank(ascending=True)
            )
            slope_df = slope_df.sort_values("Combined_rank")
            slope_df = slope_df[slope_df["p-value"] <= alpha]
            slope_df.to_csv(
                table.replace(
                    "_differential_expression.tsv", "dose_response_curves.tsv"
                ),
                sep="\t",
            )

            return slope_df

        def make_concentration_graphic(rows, xvals, direction, treatment, plot_path):
            row_df = pd.concat(
                [
                    pd.concat(rows, axis=0).reset_index(),
                    pd.concat(xvals, axis=0).reset_index(),
                    pd.concat(direction, axis=0).reset_index(),
                ],
                axis=1,
                ignore_index=True,
            )
            row_df.columns = [
                "_",
                "Normalised absolute counts",
                "_2",
                "Concentrations",
                "_3",
                "Direction interaction",
            ]
            row_df["Concentrations"] = row_df["Concentrations"].astype(str)
            fig = plt.figure(figsize=(10, 10))
            # sns.lineplot(
            try:
                sns.violinplot(
                    data=row_df,
                    x="Concentrations",
                    y="Normalised absolute counts",
                    hue="Direction interaction",
                    # err_style="bars",
                    # errorbar=("se", 2),
                    palette="Pastel1",
                    split=True,
                    cut=0,
                    gap=1,
                    # inner = "stick",
                    density_norm="area",
                )
                ax = plt.gca()
                # ax.set_xscale('log')

                fig.savefig(f"{plot_path}/{treatment}_concentration_graph.png", dpi=300)
            except ValueError:
                return ()

        def process_by_table(
            table,
            gene_count_df,
            alpha,
            cores,
            plot_path,
            model_sigmoid_func,
            meta_map,
            dmso_df,
        ):
            """
            Wrapper around model_sigmoids so that each gene table can get processed individually but in parallel


            Parameters
            ----------
            table : string
                Path of the table with Deseq results per
                treatment

            gene_count_df : object
                Pandas dataframe of gene counts per concentration

            alpha : float
                Significance level. Usually 0.05

            cores : integer
                Number of cores to use

            plot_path : string
                Path to put plots in.

            model_sigmoid_func : function
                model_sigmoid function

            meta_map : string
                Pickle file containing a dictionary with gene names as keys and annotation as values

            dmso_df : object
                Pandas dataframe containing the counts
                and concentration of the DMSO well to plot
                as a control
            """

            # Get name of treatment
            treatment = table.split("/")[-1].split("_")[0]

            # Read in corrected reads and do some data processing
            data_frame = pd.read_table(table, sep="\t", index_col=0)

            data_frame = data_frame.fillna(1)
            data_frame = data_frame.sort_values("padj")
            data_frame = data_frame[data_frame["padj"] <= alpha]

            # If there are no significant genes found abort
            if data_frame.empty:
                print(
                    f"""No significant genes found for treatment {treatment}
                    at alpha {alpha}. Skipping """
                )
                return

            # Extract only one treatment
            new_df = gene_count_df[gene_count_df["Treatments"] == treatment]
            if new_df["Concentrations"].nunique() <= 1:
                print(
                    f"Found only one concentration for treatment {treatment}. Skipping"
                )
                return

            new_df = new_df.sort_values("Concentrations")
            concentrations = new_df["Concentrations"]
            new_df = new_df[data_frame.index].fillna(0)
            new_df = new_df.loc[:, (new_df.sum() > 0)]

            if new_df.empty:
                print(
                    f"""No significant genes found for treatment {treatment}
                    at alpha {alpha}. Skipping """
                )
                return

            # Fit Sigmoid curve in aparallel over all genes
            func = partial(
                model_sigmoid_func,
                dmso_data=dmso_df,
                xvals=concentrations,
                median=list(concentrations.astype(str)),
                n_param=4,
                mode="min",
                meta_map=meta_map,
            )

            input_tuples = [
                (
                    x[0],
                    x[1],
                    dmso_df[[x[0], "Concentrations"]],
                    data_frame.loc[x[0], "padj"],
                    data_frame.loc[x[0], "log2FoldChange"],
                    treatment,
                )
                for x in new_df.items()
            ]
            input_tuples = sorted(input_tuples, key=lambda x: x[3])

            if len(input_tuples) > 100:
                input_tuples = input_tuples[:101]
            with Pool(cores) as pool:
                output_tuples = list(
                    # tqdm.tqdm(
                    map(func, input_tuples),
                    # total=input_tuples[:101],
                    # desc="Looping through genes",
                    # position=1,
                )
                # )

            # Sort the output by the goodness of fit of the fitted sigmoid
            output_tuples = sorted(output_tuples, key=lambda x: x[3], reverse=True)

            # Separate outputs
            rows, xvals, direction, r2, gene_names, slope_stats, curve_figures = zip(
                *output_tuples
            )

            # Make concentration graph
            make_concentration_graphic(rows, xvals, direction, treatment, plot_path)

            # Generate a dictionary mapping gene names to fitted curves
            figure_dict = dict(zip(gene_names, curve_figures))

            # Make a dataframe of all the fitted sigmoid stats
            slope_df = pd.DataFrame(slope_stats, index=gene_names)
            slope_df = slope_df.merge(data_frame, left_index=True, right_index=True)
            slope_df = process_slope_df(slope_df, alpha, table)

            # All the interesting sigmoid stats are given ranks to sort by
            ranks = [x for x in slope_df.columns if "rank" in x]

            # Output all top 4 in all possible stats out of
            # pvalue R2 EC50 and combined rank
            for i in range(1, 4):
                for rank in ranks:
                    gene_name = slope_df.loc[slope_df[rank] == i, "Gene"]
                    if not gene_name.empty:
                        fig = figure_dict[list(gene_name)[0]]
                        fig.savefig(
                            f"{plot_path}/{treatment}_top_{i}_in_{rank}_dose_curve",
                            dpi=200,
                        )

            # Put it all in a pdf
            pdf_file = PdfPages(f"{plot_path}/{treatment}_dose_response_curves.pdf")
            top_combined_rank = list(
                slope_df.loc[slope_df["Combined_rank"] <= 100, "Gene"]
            )

            curve_figures = [figure_dict[x] for x in top_combined_rank]
            _ = list(map(lambda x: pdf_file.savefig(x), curve_figures))
            pdf_file.close()

        # Get all different DeSeq outputs per treatment
        differential_expression_per_treatment = list(
            glob.glob(f"{self._results_path}/*_differential_expression.tsv")
        )

        # Get the summary deseq output over all gene

        gene_count_df = pd.read_csv(
            f"{self._results_path}/corrected_count_matrix.csv", index_col=0
        )

        # Get well metadta
        meta_data_df = pd.read_csv(
            f"{self._results_path}/plate_matrix.csv", index_col=0
        )

        gene_count_df.drop(columns=["Treatments", "Concentration"], inplace=True)
        # Merge all the tables
        gene_count_df = pd.merge(
            gene_count_df, meta_data_df, left_index=True, right_index=True
        )
        dmso_df = gene_count_df[gene_count_df["Treatments"] == "DMSO"]

        gene_count_df.index = gene_count_df.index.astype(str)
        # Apply the model sigmoids over all treatments. Within the treatments
        # all genes will also be handled individually

        sns.set_theme()

        func = partial(
            process_by_table,
            gene_count_df=gene_count_df,
            alpha=self._alpha,
            plot_path=self._plot_path,
            cores=self._cores,
            model_sigmoid_func=self.model_sigmoids,
            meta_map=self._meta_map,
            dmso_df=dmso_df,
        )

        _ = list(
            tqdm.tqdm(
                map(func, differential_expression_per_treatment),
                total=len(differential_expression_per_treatment),
                desc="Looping through treatments",
                position=0,
                leave=True,
            )
        )


class PowerPoint:
    """
    Class to create powerpoint presentation at the end of the BRIDGE pipeline
    """

    def __init__(self, results_path):
        """
        Intialisation of PowerPoint class

        Parameters
        ----------
        self : object
            Instance of PowerPoint class

        results_path: string
             Path of the Deseq results folder containing the tsv tables. Usually under <experiment folder>/DeSeq/Tables/

        """

        # Set all the different paths for the outputs to draw together
        # Into the pipeline

        self._results_path = results_path
        self._deseq_dir = results_path.replace("Tables", "Plots")
        self._base_dir = "/".join(results_path.split("/")[:-3]) + "/"
        self._star_dir = f"{self._base_dir}/STAR/"
        self._fastqc_dir = f"{self._base_dir}/QC/"
        self._accession = self._base_dir.split("/")[-2]

    def initialise_pres(self):
        """
        Intialisation of the presentation from a template

        Parameters
        ----------
        self : object
            Instance of PowerPoint class

        Returns
        -------
        presentation: object
            pptx object containing the presentation

        """

        # This is the powerpoint template to use
        template = "/home/ubuntu/Scripts/Headline.pptx"
        presentation = Presentation(template)

        # Get the title slide layout
        title_slide_layout = presentation.slide_layouts[0]

        # Make the title slide
        slide = presentation.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = self._accession
        subtitle.text = f"Automatically generated on: \n{date.today()}"

        return presentation

    def make_chapter(self, title_text, subtitle_text):
        """
        Make a chapter slide

        Parameters
        ----------
        self : object
            Instance of the PowerPoint class
        title_text : string
            Text to put as title

        subtitle_text :  string
            Text to put as the subtitle

        Returns
        -------
        presentation: object
            pptx object containing the presentation

        """

        presentation = self.presentation

        title_slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = title_text
        subtitle.text = subtitle_text

        return presentation

    def make_fastqc_slides(self, fastqc_zip):
        """
        Make slides with the fastqc outputs

        Parameters
        ----------
        self : object
            Instance of PowerPoint class

        fastqc_zip : generator
            zipped generator containing all different fastqc zips and
            their respective titles

        """
        presentation = self.presentation

        # Get name of the experiment, plate info and read info
        folder_name = fastqc_zip.split("/")[-1].split(".")[0]
        try:
            plate_info = folder_name.split("_")[0]
            read_info = folder_name.split("_")[-2]
        except:
            plate_info = fastqc_zip.split("/")[-1]
            plate_info = fastqc_zip.split("/")[-1]
        plate_info = plate_info.replace("plate", "")

        # Define which plots we're interested in
        plots_of_interest = [
            "adapter_content.png",
            "per_base_sequence_content.png",
            "per_sequence_quality.png",
            "per_sequence_gc_content.png",
        ]

        # Open the fastqc zip to extract the above plots
        with ZipFile(fastqc_zip) as zip_inf:
            list(
                map(
                    lambda x: zip_inf.extract(f"{folder_name}/Images/{x}"),
                    plots_of_interest,
                )
            )

        # Make the slides
        slide_layout = presentation.slide_layouts[10]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"FASTQC Plate: {plate_info} Read: {read_info}"

        # Place images
        left = Cm(5)
        top = Cm(4)
        right = Cm(20)
        bottom = Cm(12)
        slide.shapes.add_picture(
            f"{folder_name}/Images/per_sequence_quality.png",
            left,
            top,
            width=Cm(10),
            height=Cm(7),
        )
        slide.shapes.add_picture(
            f"{folder_name}/Images/per_base_sequence_content.png",
            right,
            top,
            width=Cm(10),
            height=Cm(7),
        )
        slide.shapes.add_picture(
            f"{folder_name}/Images/adapter_content.png",
            left,
            bottom,
            width=Cm(10),
            height=Cm(7),
        )
        slide.shapes.add_picture(
            f"{folder_name}/Images/per_sequence_gc_content.png",
            right,
            bottom,
            width=Cm(10),
            height=Cm(7),
        )

        self.presentation = presentation

    def make_star_slides(self, star_file):
        """
        Make slides containing the STARsolo outputs

        Parameters
        ----------
        self : object
            Instance of PowerPoint class

        star_file : string
            Path of the Log.final.out file generated by STARsolo

        """
        star_file, extra_text = star_file

        presentation = self.presentation
        if "summary_blank" in star_file:
            plate_info = "0"
        elif "summary" in star_file:
            plate_info = star_file.split("/")[-2].split("_")[0]
        else:
            plate_info = (
                star_file.split("/")[-1]
                .split("_")[0]
                .split(".")[0]
                .replace("plate", "")
            )

        slide_layout = presentation.slide_layouts[10]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"STARsolo Plate: {plate_info}"
        top = width = height = Cm(2)
        left = Cm(10)
        txt_box = slide.shapes.add_textbox(left, top, width, height)
        txt_frame = txt_box.text_frame
        par = txt_frame.add_paragraph()
        par.font.size = Pt(10)

        if star_file == "summary_blank.csv":
            par.text = extra_text
        else:
            with open(star_file) as inf:
                par.text = inf.read() + extra_text

        self.presentation = presentation

    def calculate_sum_counts(self, mtx):
        return pd.DataFrame.sparse.from_spmatrix(mmread(mtx)).to_numpy().sum()

    def make_ratio_deduplication(self, input_tuple):
        mtx_deduplicated, mtx_not_deduplicated = input_tuple
        counts_dedup = self.calculate_sum_counts(mtx_deduplicated)
        counts_no_dedup = self.calculate_sum_counts(mtx_not_deduplicated)
        ratio = 1 - (counts_dedup / counts_no_dedup)
        text = f"""\n\n\n
            Reads before de-duplication: {counts_no_dedup}
            Reads after de-duplication: {counts_dedup}
            Fraction non-unique reads (these are thrown away): {ratio:.6f}"""
        return text

    def parse_dose_response_plot_title(self, file_name):
        """
        Parse name for dose response plots

        Parameters
        ----------
        file_name : string
            Name of the file

        Returns
        -------
        string
            Parse plot title

        """

        file_name = file_name.split("/")[-1].replace("_dose_curve.png", "")
        treatment = file_name.split("_")[0]
        description = " ".join(file_name.split("_")[1:])

        return f"{treatment}: 6. Dose response plot {description}"

    def make_deseq_slides(self):
        """
        Make slides containing the DeSeq plots.

        Parameters
        ----------
        self : object
            Instance of PowerPoint class

        """
        presentation = self.presentation

        # List of plots to include
        important_plots = [
            "ReadsPerWell.png",
            "replicate_correlations.png",
            "read_decay_curves.png",
            "Raw_counts_distribution_plot.png",
            "Corrected_counts_distribution_plot.png",
            "Plate_uncorrected_pca.png",
            "Plate_corrected_treatments_pca.png",
            "Plate_corrected_readsums_pca.png",
            "Outlier_removed_pca.png",
            "2D_umap.png",
            "clusters_membership.png",
            "correlation_treatments.png",
            "correlation_clusters.png",
        ]

        important_plots_broad = [
            "louvain_clustering_umap.png",
            "louvain_clustering_pca.png",
            "leiden_clustering_umap.png",
            "leiden_clustering_pca.png",
            "per_cluster_logfold.png",
            "per_cluster_expression.png",
        ]

        # Titles for these plots
        import_plot_titles = [
            "Number of aligned reads per well",
            "Correlations between replicates",
            "Read decay curves",
            "Distribution of raw read counts mapped to reads",
            """Distribution of corrected read counts mapped to reads""",
            """Principal component analysis of wells
            before plate batch coffection and outlier removal""",
            """Principal component analysis of wells after
            plate batch coffection and before outlier removal""",
            """Principal component analysis of wells with
            colouring by readsums after plate batch coffection
            and before outlier removal""",
            """Principal component analysis of wells
            after plate batch coffection and outlier removal""",
            "Normal 2D PCA",
            "Cluster membership LEIDEN heatmap",
            "Correlation normalised gene expression per treatment",
            "Correlation normalised gene expression per cluster",
        ]

        import_plot_titles_broad = [
            "UMAP of Louvain clustering",
            "PCA of Louvain clustering",
            "UMAP of Leiden clustering",
            "PCA of Leiden clustering",
            "Logfold change of gene expression per LEIDEN cluster",
            "Normalised gene expression per LEIDEN cluster",
        ]

        # Generate the slides
        important_plots = [f"{self._deseq_dir}{x}" for x in important_plots]
        important_plot_zip = zip(important_plots, import_plot_titles)
        list(map(self.make_picture_slide, important_plot_zip))

        important_plots_broad = [f"{self._deseq_dir}{x}" for x in important_plots_broad]
        important_plot_broad_zip = zip(important_plots_broad, import_plot_titles_broad)
        list(map(self.make_picture_slide_broad, important_plot_broad_zip))

        self.presentation = self.make_chapter(
            "Functional analysis", "Slides ordered by treatment"
        )

        # Generate volcano plots
        volcano_plots = list(glob.glob(f"{self._deseq_dir}*_volcano_plot.png"))
        volcano_titles = [
            f"""{x.split("/")[-1].split("_")[0]}: 1. Volcano plot"""
            for x in volcano_plots
        ]
        plot_zips = list(zip(volcano_plots, volcano_titles))

        slide_functions = [self.make_picture_slide_volcano] * len(volcano_plots)

        broad_format_substrings = ["bubble", "circle", "cluster"]
        go_plots = list(glob.glob(f"{self._deseq_dir}*_GO_*.png"))
        go_plots_broad = [
            x for x in go_plots if any(y in x for y in broad_format_substrings)
        ]
        go_plots_slim = list(set(go_plots) - set(go_plots_broad))

        functional_slides_titles_slim = [
            f"""{x.split("/")[-1].split("_")[0]}: 2. Gene ontology analysis of {x.split("/")[-1].split("_")[2]}-regulated genes: 
            """
            for x in go_plots_slim
        ]
        functional_slides_titles_broad = [
            f"""{x.split("/")[-1].split("_")[0]}: 3. Gene ontology analysis of both up and down-regulated genes: 
            """
            for x in go_plots_broad
        ]

        plot_zips += list(zip(go_plots_slim, functional_slides_titles_slim))
        plot_zips += list(zip(go_plots_broad, functional_slides_titles_broad))
        slide_functions += [self.make_picture_slide] * len(go_plots_slim)
        slide_functions += [self.make_picture_slide_broad] * len(go_plots_broad)

        kegg_plots = list(glob.glob(f"{self._deseq_dir}*_KEGG_*.png"))
        functional_slides_titles = [
            f"""{x.split("/")[-1].split(".")[1]}: 4. KEGG enrichment analysis overview
            """
            for x in kegg_plots
        ]
        plot_zips += list(zip(kegg_plots, functional_slides_titles))
        slide_functions += [self.make_picture_slide] * len(kegg_plots)

        kegg_plots = list(glob.glob(f"{self._deseq_dir}ko*.*.png"))
        functional_slides_titles = [
            f"""{x.split("/")[-1].split(".")[1]}: 5. KEGG enrichment analysis individual pathways
            """
            for x in kegg_plots
        ]
        plot_zips += list(zip(kegg_plots, functional_slides_titles))

        slide_functions += [self.make_picture_slide_broad] * len(kegg_plots)

        # Generate gene atlas slides
        # gene_atlas_plots = list(
        #    glob.glob(f"{self._deseq_dir}*_gene_atlas.png"))
        # gene_atlas_titles = [
        #    f"""{x.split("/")[-1].split("_")[0]}: 6. Transcriptomic fingerprint on UMAP of differential expression""" for x in gene_atlas_plots]
        # plot_zips += list(zip(gene_atlas_plots, gene_atlas_titles))
        # slide_functions += [self.make_picture_slide] * len(gene_atlas_plots)

        # Generate dose response slides
        dose_response_plots = list(
            glob.glob(f"{self._deseq_dir}*Combined_rank*_dose_curve.png")
        )
        dose_response_titles = [
            self.parse_dose_response_plot_title(x) for x in dose_response_plots
        ]
        plot_zips += list(zip(dose_response_plots, dose_response_titles))
        slide_functions += [self.make_picture_slide_broad] * len(dose_response_plots)

        try:
            _, titles = zip(*plot_zips)
        except ValueError:
            print(f"No DeSeq plots found in {self._deseq_dir}")
            return

        function2plot_zip = list(zip(slide_functions, plot_zips, titles))
        function2plot_zip = sorted(function2plot_zip, key=lambda x: x[-1])
        old_treatment = ""
        for func, arg, titles in function2plot_zip:
            treatment = titles.split(":")[0]
            if treatment != old_treatment:
                self.presentation = self.make_chapter(treatment, "Functional analysis")
                old_treatment = treatment

            try:
                func(arg)
            except:
                print(f"Could not make slide for {arg}")

        self.presentation = presentation

    def make_picture_slide_dr(self, input_tuple):
        """
        Make dose response slides

        Parameters
        ----------
        self : object
            Instance of PowerPoint class

        input_tuple : tuple
            Tuple containing the plot to include and its title

        """

        image, slide_title = input_tuple
        presentation = self.presentation
        slide_layout = presentation.slide_layouts[10]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"DeSeq analysis: {slide_title}"
        left = Cm(1)

        top = Cm(3)

        slide.shapes.add_picture(image, left, top, width=Cm(30), height=Cm(15))
        self.presentation = presentation

    def make_picture_slide_volcano(self, input_tuple):
        """
        Make volcano plotslides

        Parameters
        ----------
        self : object
            Instance of PowerPoint class

        input_tuple : tuple
            Tuple containing the plot to include and its title

        """

        image, slide_title = input_tuple
        image = image.replace("_A5_", "_")
        hyperlinks = image.replace("_volcano_plot.png", "_uniprot_links.txt")

        presentation = self.presentation
        slide_layout = presentation.slide_layouts[8]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = ""
        title.text = slide_title

        left = Cm(1.5)
        top = Cm(3.5)
        width = Cm(15)
        height = Cm(15)
        slide.shapes.add_picture(image, left, top, width=width, height=height)

        with open(hyperlinks, "r") as inf:
            hyperlinks = inf.read().split("\n")
        text_frame = slide.shapes[3].text_frame
        text_frame.clear()  # not necessary for newly-created shape
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = (
            "Please click on the following to go to the respective Uniprot page\n"
        )
        font = run.font
        font.name = "Calibri"
        font.size = Pt(12)

        for i, hit in enumerate(hyperlinks[:10]):
            gene, hyperlink = hit.split("\t")
            run = p.add_run()
            run.text = f"{i+1}: {gene}\n"
            font = run.font
            font.name = "Calibri"
            font.size = Pt(12)
            font.italic = None
            run.hyperlink.address = hyperlink
        self.presentation = presentation

    def make_picture_slide(self, input_tuple):
        """
        Make general slides

        Parameters
        ----------
        self : object
            Instance of PowerPoint class

        input_tuple : tuple
            Tuple containing the plot to include and its title

        """

        image, slide_title = input_tuple
        if not os.path.isfile(image):
            return
        presentation = self.presentation
        slide_layout = presentation.slide_layouts[10]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = slide_title
        left = Cm(7.5)
        top = Cm(3)
        width = Cm(15.5)
        height = Cm(15.5)
        slide.shapes.add_picture(image, left, top, height=height)
        self.presentation = presentation

    def make_picture_slide_broad(self, input_tuple):
        """
        Make general slides

        Parameters
        ----------
        self : object
            Instance of PowerPoint class

        input_tuple : tuple
            Tuple containing the plot to include and its title

        """

        image, slide_title = input_tuple
        if not os.path.isfile(image):
            return
        presentation = self.presentation
        slide_layout = presentation.slide_layouts[10]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = slide_title
        left = Cm(3.5)
        top = Cm(3.5)
        width = Cm(30)
        height = Cm(15)
        slide.shapes.add_picture(image, left, top, width=width)
        self.presentation = presentation

    def make_presentation(self):
        """
        Wrapper function around making the output presentation

        Parameters
        ----------
        self : object
            Instance of PowerPoint class

        """
        # Initialise presentation
        presentation = self.initialise_pres()
        self.presentation = presentation

        # Fastqc chapter
        fastq_zip = list(glob.glob(self._fastqc_dir + "*.zip"))
        fastq_zip = sorted(fastq_zip)
        self.presentation = self.make_chapter(
            "FastQC analysis", "Checking read quality"
        )
        list(map(self.make_fastqc_slides, fastq_zip))

        # Starsolo chapter
        self.presentation = self.make_chapter("STARsolo analysis", "Aligning to genome")
        star_outs = glob.glob(self._star_dir + "*.out")
        summary_mapping = glob.glob(self._star_dir + "*_CountTables/summary.csv")
        print(self._star_dir + "*_CountTables/umiDedup-NoDedup.mtx")

        matrix_no_deduplicated = glob.glob(
            self._star_dir + "*_CountTables/umiDedup-NoDedup.mtx"
        )
        matrix_deduplicated = glob.glob(
            self._star_dir + "*_CountTables/umiDedup-1MM_All.mtx"
        )

        star_outs = list(zip(star_outs, [""] * len(star_outs)))
        ratios_deduplication = list(
            map(
                self.make_ratio_deduplication,
                zip(matrix_deduplicated, matrix_no_deduplicated),
            )
        )
        if len(summary_mapping) > 0:
            summary_mapping = list(zip(summary_mapping, ratios_deduplication))
        else:
            summary_mapping = list(
                zip(
                    ["summary_blank.csv"] * len(ratios_deduplication),
                    ratios_deduplication,
                )
            )

        list(map(self.make_star_slides, star_outs))

        list(map(self.make_star_slides, summary_mapping))

        # DeSeq chapter
        self.presentation = self.make_chapter(
            "DeSeq analysis", "Differential expression"
        )
        # try:
        self.make_deseq_slides()
        # except ValueError:
        #    print("No DeSeq results")
        self.presentation = self.make_chapter("That is all", "Thank you")
        pres_loc = f"{self._base_dir}/{self._accession}_BRIDGE_summary.pptx"
        print(f"Dropping Presetation in: {pres_loc}")
        # Save presentations
        self.presentation.save(pres_loc)
